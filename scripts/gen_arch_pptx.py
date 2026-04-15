#!/usr/bin/env python3
"""
Generate an editable PowerPoint architecture diagram for the
Adaptive RAG Clinical Assistant system.

Content verified against the actual codebase (src/config.py, src/retrieval.py,
src/personas.py, src/query_classifier.py, src/llm/).

Usage:
    python scripts/gen_arch_pptx.py

Output:
    docs/architecture_diagram.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent.parent / "docs" / "architecture_diagram.pptx"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY = "#1A1A2E"
GRAY = "#6B7280"
D_GRAY = "#374151"
WHITE = "#FFFFFF"
BG = "#F8FAFC"

# ── Layer definitions (content verified against codebase) ────────────────────
# Verified sources:
#   Chunk size/overlap   → src/config.py (DEFAULT_CHUNK_SIZE=800, DEFAULT_CHUNK_OVERLAP=150)
#   Embedding models     → src/config.py (EMBEDDING_MODELS dict)
#   RRF k constant       → src/config.py (RRF_K_CONSTANT=60)
#   Fusion algorithm     → src/retrieval.py (ReciprocalRankFusion, NOT weighted average)
#   Persona types        → src/personas.py (5 UserType enum values)
#   Query types          → src/query_classifier.py (9 QueryType enum values)
#   LLM models           → src/config.py (OPENAI_MODELS, OLLAMA_MODELS)
LAYERS = [
    {
        "title": "USER INTERFACE",
        "subtitle": "app.py  ·  Streamlit web application",
        "accent": "#1D4ED8",
        "fill": "#DBEAFE",
        "cards": [
            ("User Profile", "Role + Experience Years\n→ detect_user_type()"),
            ("Document Upload", "PDF Clinical Documents\n(PyPDFLoader)"),
            ("Query Interface", "Natural Language\nQuestion Input"),
            ("Chat History", "Session Cache\nFormatted Responses"),
        ],
    },
    {
        "title": "USER & QUERY PROFILING",
        "subtitle": "src/personas.py  ·  src/query_classifier.py",
        "accent": "#7C3AED",
        "fill": "#EDE9FE",
        "cards": [
            (
                "Persona Detection",
                "NOVICE · INTERMEDIATE\nEXPERT · REGULATORY\nEXECUTIVE\nrule-based, deterministic",
            ),
            (
                "Query Classification",
                "9 types: Definition\nProcedure · Compliance\nSafety · Eligibility\nComparison + more",
            ),
            (
                "Response Config",
                "detail_level · max_length\nuse_tables · use_bullets\ninclude_references\nformatting hints",
            ),
        ],
    },
    {
        "title": "HYBRID RETRIEVAL PIPELINE",
        "subtitle": (
            "src/retrieval.py  ·  HybridRetriever  ·  "
            "ReciprocalRankFusion  (k=60, Cormack et al. 2009)"
        ),
        "accent": "#A21CAF",
        "fill": "#FDF4FF",
        "cards": [
            (
                "Semantic Search\n(ChromaDB)",
                "HuggingFace Sentence Transformers\nDefault: S-PubMedBert-MS-MARCO\n768 dim · cosine similarity",
            ),
            (
                "Lexical Search\n(BM25)",
                "rank-bm25 library\nTF-IDF term-frequency\nExact keyword matching",
            ),
            (
                "RRF Fusion  ( k = 60 )",
                "score(d) = Σ 1/(k + rank)\nrank-based merge\nTop-5 deduplicated docs",
            ),
        ],
    },
    {
        "title": "VECTOR STORE  (ChromaDB)",
        "subtitle": (
            "Persisted: ./chroma_adaptive_{model}/    "
            "chunk_size = 800 tokens    chunk_overlap = 150"
        ),
        "accent": "#D97706",
        "fill": "#FEF3C7",
        "cards": [
            ("IRC Charters", "Decision criteria\n& committee roles"),
            ("Clinical Protocols", "Study design\n& procedures"),
            ("RECIST 1.1\nGuidelines", "Tumor response\ncriteria"),
            ("PET/CT Protocols\n& Vendor Docs", "Imaging instructions\n& equipment specs"),
        ],
    },
    {
        "title": "RESPONSE GENERATION",
        "subtitle": (
            "src/prompts.py  ·  src/llm/factory.py  ·  "
            "ResponseStyler  ·  LLMFactory.create(config)"
        ),
        "accent": "#059669",
        "fill": "#ECFDF5",
        "cards": [
            (
                "Adaptive Prompt\nBuilder",
                "query + retrieved chunks\n+ ResponseConfig →\npersona-aware prompt",
            ),
            (
                "OpenAI  (Cloud)",
                "gpt-4o-mini  ·  gpt-4o\ngpt-4  ·  gpt-3.5-turbo\nOpenAI REST API",
            ),
            (
                "Ollama  (Local /\nAir-gapped)",
                "llama3.1:8b  ·  BioMistral 7B\nmistral:7b  ·  MedGemma 4B\nfully on-premise",
            ),
        ],
    },
    {
        "title": "RESPONSE WITH INLINE CITATIONS",
        "subtitle": (
            "Persona-tailored answer  ·  Citation-backed  ·  "
            "Structured per ResponseConfig  ·  Detail level adapted to user expertise"
        ),
        "accent": "#0284C7",
        "fill": "#E0F2FE",
        "cards": [],
    },
]

# ── Layout constants (inches) ─────────────────────────────────────────────────
SLIDE_W = 13.33
SLIDE_H = 7.5
BX = 0.15  # band left edge
BW = 13.03  # band width
BH = 1.03  # standard band height (bands with cards)
BH_LAST = 0.62  # last band height (output, no cards)
ARROW_H = 0.14  # arrow height between bands
Y0 = 0.44  # y-start of first band

# Within each band:
HEADER_H = 0.26  # colored header bar height
SUBTITLE_H = 0.14  # subtitle row height
CARD_TOP_OFFSET = HEADER_H + SUBTITLE_H + 0.07  # card y relative to band top
L_PAD = 0.1  # left padding inside band


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _rect(slide, x, y, w, h, fill_hex, border_hex=None, border_pt=0.75):
    """Add a plain rectangle shape."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = _rgb(border_hex)
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def _rounded_rect(slide, x, y, w, h, fill_hex, border_hex=None, border_pt=0.75):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = _rgb(border_hex)
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def _textbox(slide, x, y, w, h, text, size, bold=False, italic=False,
             color=NAVY, align=PP_ALIGN.LEFT, wrap=False):
    """Add a simple single-paragraph textbox."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return txb


def _draw_band(slide, by: float, layer: dict, bh: float) -> None:
    """Draw one full-width architecture layer band."""
    accent = layer["accent"]
    fill = layer["fill"]
    cards = layer["cards"]
    content_w = BW - L_PAD * 2  # usable width inside the band

    # 1. Band background (light fill + accent border)
    _rect(slide, BX, by, BW, bh, fill, accent, 1.0)

    # 2. Solid accent header bar
    _rect(slide, BX, by, BW, HEADER_H, accent)

    # 3. Layer title text (white bold, inside header)
    _textbox(slide, BX + L_PAD, by + 0.04, content_w, HEADER_H - 0.05,
             layer["title"], size=10, bold=True, color=WHITE)

    # 4. Subtitle (gray italic, below header)
    _textbox(slide, BX + L_PAD, by + HEADER_H + 0.01, content_w, SUBTITLE_H,
             layer["subtitle"], size=7, italic=True, color=D_GRAY)

    if not cards:
        return

    # 5. Cards row
    n = len(cards)
    gap = 0.1
    card_w = (content_w - gap * (n - 1)) / n
    card_h = bh - CARD_TOP_OFFSET - 0.07
    card_y = by + CARD_TOP_OFFSET

    for i, (title, subtitle) in enumerate(cards):
        cx = BX + L_PAD + i * (card_w + gap)
        card = _rounded_rect(slide, cx, card_y, card_w, card_h, WHITE, accent, 0.75)

        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Card title (bold)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = _rgb(NAVY)

        # Card subtitle lines
        for line in subtitle.split("\n"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = line
            run2.font.size = Pt(6.5)
            run2.font.color.rgb = _rgb(GRAY)


def _draw_arrow(slide, cx: float, y: float, color: str) -> None:
    """Draw a small ▼ arrow between bands."""
    txb = slide.shapes.add_textbox(
        Inches(cx - 0.15), Inches(y + 0.01), Inches(0.3), Inches(0.12)
    )
    p = txb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "▼"
    run.font.size = Pt(9)
    run.font.color.rgb = _rgb(color)


def build_diagram() -> None:
    """Build and save the editable PowerPoint architecture diagram."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Blank slide (layout index 6 is typically the blank layout)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Slide background
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)

    # Slide title
    _textbox(
        slide, 0.3, 0.07, 12.73, 0.32,
        "Adaptive RAG Clinical Assistant  —  System Architecture",
        size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )

    y = Y0
    for i, layer in enumerate(LAYERS):
        is_last = i == len(LAYERS) - 1
        bh = BH_LAST if is_last else BH
        _draw_band(slide, y, layer, bh)
        y += bh
        if not is_last:
            _draw_arrow(slide, BX + BW / 2, y, color=LAYERS[i + 1]["accent"])
            y += ARROW_H

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build_diagram()
